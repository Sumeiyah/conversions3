import os
import subprocess
from shutil import which
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import io

UPLOAD_FOLDER = 'uploads/'

def convert_file(conversion_type, filepath):
    filename, file_extension = os.path.splitext(filepath)
    output_extension = get_output_extension(conversion_type)
    output_filename = f"{filename}.{output_extension}"
    output_filepath = os.path.join(UPLOAD_FOLDER, os.path.basename(output_filename))

    if conversion_type == 'pdf_to_word':
        if not can_convert_pdf_to_word():
            return None, None
        convert_pdf_to_word(filepath, output_filepath)
    elif conversion_type == 'word_to_jpg':
        convert_word_to_jpg(filepath, output_filepath)
    elif conversion_type == 'word_to_pdf':
        if not can_convert_office_to_pdf():
            return None, None
        convert_word_to_pdf(filepath, output_filepath)
    elif conversion_type == 'ppt_to_pdf':
        if not can_convert_office_to_pdf():
            return None, None
        convert_ppt_to_pdf(filepath, output_filepath)
    elif conversion_type == 'pdf_to_ppt':
        convert_pdf_to_ppt(filepath, output_filepath)
    else:
        raise ValueError("Unsupported conversion type")

    return output_filepath, output_extension

def can_convert_pdf_to_word():
    return which('pdftoppm') and which('pdfinfo')

def can_convert_office_to_pdf():
    return which('libreoffice')

def convert_word_to_jpg(input_path, output_path):
    document = Document(input_path)
    img_path = f"{output_path}.jpg"
    document.save(img_path)
    return img_path

def convert_pdf_to_word(input_path, output_path):
    doc = Document()
    images = convert_from_path(input_path)
    for img in images:
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        doc.add_picture(img_bytes)
    doc.save(output_path)

def convert_word_to_pdf(input_path, output_path):
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', input_path, '--outdir', UPLOAD_FOLDER])
    return output_path

def convert_ppt_to_pdf(input_path, output_path):
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', input_path, '--outdir', UPLOAD_FOLDER])
    return output_path

def convert_pdf_to_ppt(input_path, output_path):
    prs = Presentation()
    slides = convert_from_path(input_path)
    for slide in slides:
        slide_layout = prs.slide_layouts[5]
        slide_obj = prs.slides.add_slide(slide_layout)
        img_path = '/tmp/slide.jpg'
        slide.save(img_path, 'JPEG')
        slide_obj.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)
    prs.save(output_path)

def get_output_extension(conversion_type, pure_extension=False):
    extensions = {
        'pdf_to_word': 'docx',
        'pdf_to_jpg': 'jpg',
        'word_to_pdf': 'pdf',
        'ppt_to_pdf': 'pdf',
        'pdf_to_ppt': 'pptx',
        'word_to_jpg': 'jpg'
    }
    return extensions.get(conversion_type, '').lstrip('.') if pure_extension else extensions.get(conversion_type, '')
